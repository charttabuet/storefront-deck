# -*- coding: utf-8 -*-
"""
deck_core.py  ---  แกนกลางสร้างสไลด์ (ธีมขาว-แดง Tefal)
build_deck(branches, campaign, month, template_path, out_path)
  branches = [ {"name": "สาขา...", "note": "...", "photos": [path1, path2, ...]}, ... ]

*** ปรับสีแบรนด์ได้ที่ตัวแปรด้านล่างนี้ที่เดียว ***
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ================== ธีมสี (แก้ตรงนี้) ==================
RED       = RGBColor(0xE2, 0x00, 0x1A)   # แดง Tefal (สีหลัก)
RED_DK    = RGBColor(0xB3, 0x00, 0x15)   # แดงเข้ม (accent)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x22, 0x22, 0x22)   # ข้อความเข้ม
GREY      = RGBColor(0x6B, 0x72, 0x80)   # ข้อความจาง / ฟุตเตอร์
LIGHTRED  = RGBColor(0xFF, 0xD6, 0xDA)   # ชมพูอ่อน (ป้ายบนแถบแดง)
BORDER    = RGBColor(0xDD, 0xDD, 0xDD)   # ขอบรูป
MISS      = RGBColor(0x9A, 0x9A, 0x9A)   # เทา (สไลด์แจ้งไม่มีรูป)
FONT      = "Tahoma"
# =====================================================

MAX_PER_SLIDE = 6
SLIDE_W, SLIDE_H = 13.333, 7.5
BAND_H   = 1.05
AREA_L, AREA_T = 0.5, BAND_H + 0.30
AREA_W = SLIDE_W - 2 * AREA_L
AREA_H = SLIDE_H - AREA_T - 0.55
PAD    = 0.12


def _grid_for(n):
    return {1: (1, 1), 2: (2, 1), 3: (3, 1),
            4: (2, 2), 5: (3, 2), 6: (3, 2)}[n]


def _title_size(text, base=26, min_size=14):
    """ย่อขนาดฟอนต์หัวข้ออัตโนมัติ เผื่อชื่อยาว (ลูกค้า+สาขา+กิจกรรม)
    ให้พอดีในแถบสีบนสไลด์ ไม่ล้นออกนอกกรอบ"""
    n = len(text or "")
    if n <= 26:
        return base
    if n <= 36:
        return 22
    if n <= 48:
        return 18
    if n <= 64:
        return 16
    return min_size


def _text(slide, l, t, w, h, text, size, color, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.name = FONT; r.font.color.rgb = color
    return tb


def _rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def _place_image(slide, path, box_l, box_t, box_w, box_h):
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 4, 3
    inner_w, inner_h = box_w - PAD, box_h - PAD
    img_ar, box_ar = iw / ih, inner_w / inner_h
    if img_ar >= box_ar:
        w = inner_w; h = inner_w / img_ar
    else:
        h = inner_h; w = inner_h * img_ar
    l = box_l + (box_w - w) / 2
    t = box_t + (box_h - h) / 2
    pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1)
    return pic


def _footer(slide, campaign, month, page_no):
    _text(slide, 0.5, SLIDE_H - 0.5, 9, 0.35, f"{campaign} • {month}", 10, GREY)
    _text(slide, SLIDE_W - 2.0, SLIDE_H - 0.5, 1.5, 0.35, str(page_no), 10, GREY,
          align=PP_ALIGN.RIGHT)


def _cover_slide(prs, layout, campaign, month, kicker=None):
    slide = prs.slides.add_slide(layout)
    # พื้นปกแดงเต็มจอ
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RED)
    # บล็อกแดงเข้มมุมขวาบน เป็น motif
    _rect(slide, 10.6, 0, 2.733, 2.6, RED_DK)
    # หัวเรื่องเล็ก (ส่ง kicker มาจาก payload ได้ ถ้าไม่ส่งใช้ค่าเดิม)
    _text(slide, 0.9, 2.0, 11, 0.6,
          kicker or "รายงานภาพหน้าร้าน • CAMPAIGN REPORT",
          18, WHITE, bold=True)
    # ชื่อแคมเปญ
    _text(slide, 0.9, 2.7, 11.5, 1.8, campaign, 46, WHITE, bold=True,
          anchor=MSO_ANCHOR.TOP)
    # เดือน
    _text(slide, 0.9, 4.6, 11, 0.8, month, 26, LIGHTRED)
    # ฟุตเตอร์
    _text(slide, 0.9, 6.6, 11, 0.5, "จัดทำโดยทีม Sales • สร้างอัตโนมัติด้วยระบบ",
          12, LIGHTRED)
    return slide


def _branch_slide(prs, layout, branch, imgs, page_idx, page_total, note,
                  campaign, month, page_no):
    slide = prs.slides.add_slide(layout)
    _rect(slide, 0, 0, SLIDE_W, BAND_H, RED)
    title = branch + (f"  ({page_idx}/{page_total})" if page_total > 1 else "")
    _text(slide, 0.5, 0.05, 9.8, BAND_H - 0.1, title, _title_size(title),
          WHITE, bold=True)
    _text(slide, SLIDE_W - 3.3, 0.05, 2.8, BAND_H - 0.1, f"{len(imgs)} รูป",
          14, LIGHTRED, bold=True, align=PP_ALIGN.RIGHT)

    area_t, area_h = AREA_T, AREA_H
    if note and page_idx == 1:
        _text(slide, 0.5, BAND_H + 0.05, AREA_W, 0.3, "หมายเหตุ: " + note, 12, GREY)
        area_t += 0.28; area_h -= 0.28

    n = len(imgs)
    cols, rows = _grid_for(n)
    cell_w, cell_h = AREA_W / cols, area_h / rows
    for i, path in enumerate(imgs):
        r, c = divmod(i, cols)
        in_row = min(cols, n - r * cols)
        row_offset = (AREA_W - in_row * cell_w) / 2
        _place_image(slide, path, AREA_L + row_offset + c * cell_w,
                     area_t + r * cell_h, cell_w, cell_h)
    _footer(slide, campaign, month, page_no)


def _missing_slide(prs, layout, branch, campaign, month, page_no):
    slide = prs.slides.add_slide(layout)
    _rect(slide, 0, 0, SLIDE_W, BAND_H, MISS)
    _text(slide, 0.5, 0.05, 12.3, BAND_H - 0.1, branch, _title_size(branch),
          WHITE, bold=True)
    _text(slide, 0.5, 3.1, AREA_W, 1.2, "⚠  ยังไม่มีรูปส่งเข้ามาสำหรับสาขานี้",
          24, MISS, bold=True, align=PP_ALIGN.CENTER)
    _footer(slide, campaign, month, page_no)


def build_deck(branches, campaign, month, template_path, out_path, kicker=None):
    """สร้าง .pptx จากรายการสาขา  คืนค่า (จำนวนรูปรวม, [สาขาที่ไม่มีรูป])
    (สร้างสไลด์ปกในโค้ดเอง ไม่ต้องพึ่ง template.pptx แล้ว)"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]   # blank

    _cover_slide(prs, layout, campaign, month, kicker)

    page_no, total, missing = 1, 0, []
    for b in branches:
        name = b["name"]
        imgs = b.get("photos", [])
        note = b.get("note", "")
        if not imgs:
            page_no += 1
            _missing_slide(prs, layout, name, campaign, month, page_no)
            missing.append(name)
            continue
        total += len(imgs)
        pages = [imgs[i:i + MAX_PER_SLIDE] for i in range(0, len(imgs), MAX_PER_SLIDE)]
        for pi, chunk in enumerate(pages, start=1):
            page_no += 1
            _branch_slide(prs, layout, name, chunk, pi, len(pages), note,
                          campaign, month, page_no)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    return total, missing
